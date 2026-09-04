"""Tests unitaires pour le parseur enrichi PowerPoint (.pptx).

Vérifie l'extraction des diapositives, tableaux, images, listes imbriquées,
et des notes secrètes de l'orateur (Speaker Notes) pour la création de fiches Anki.
"""

import io
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from PIL import Image
from pptx import Presentation

from ankiforge.services.parsing.chunking_service import ChunkingService
from ankiforge.services.parsing.document_parser import DocumentParser
from ankiforge.services.parsing.pptx_parser import PptxParser


def create_sample_pptx(
    file_path: Path,
    include_notes: bool = True,
    include_table: bool = True,
    include_image: bool = True,
) -> Path:
    """Génère une présentation PowerPoint de test complète avec titre, contenu, tableau, image et notes."""
    prs = Presentation()

    # Slide 1 : Titre + Puces à plusieurs niveaux + Notes d'orateur
    slide1 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide1.shapes.title.text = "Physiopathologie Cardiaque"

    body_shape = slide1.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "Débit Cardiaque (DC) :"

    p2 = tf.add_paragraph()
    p2.text = "Formule : DC = FC x VES"
    p2.level = 1

    p3 = tf.add_paragraph()
    p3.text = "Régulation par le système nerveux autonome"
    p3.level = 1

    if include_notes:
        notes1 = slide1.notes_slide
        notes1.notes_text_frame.text = "Insister sur la loi de Frank-Starling !\nQuestion classique d'internat."

    # Slide 2 : Tableau + Image
    blank_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank_layout)

    # Titre manuel sous forme de zone de texte
    title_box = slide2.shapes.add_textbox(50, 20, 400, 50)
    title_box.text_frame.text = "Tableau des Valeurs Physiologiques"

    if include_table:
        table_shape = slide2.shapes.add_table(3, 2, 50, 100, 300, 150)
        tbl = table_shape.table
        tbl.cell(0, 0).text = "Paramètre"
        tbl.cell(0, 1).text = "Norme"
        tbl.cell(1, 0).text = "Pression Artérielle"
        tbl.cell(1, 1).text = "120/80 mmHg"
        tbl.cell(2, 0).text = "Fréquence Cardiaque"
        tbl.cell(2, 1).text = "60-100 bpm"

    if include_image:
        img_buf = io.BytesIO()
        Image.new("RGB", (30, 30), color="blue").save(img_buf, format="PNG")
        img_buf.seek(0)
        slide2.shapes.add_picture(img_buf, 400, 100)

    if include_notes:
        notes2 = slide2.notes_slide
        notes2.notes_text_frame.text = "Attention : retenir les unités (mmHg vs kPa)."

    prs.save(str(file_path))
    return file_path


def test_pptx_parser_speaker_notes(tmp_path):
    """Vérifie que les notes de l'orateur sont extraites sous forme de bloc citation dédié."""
    pptx_path = tmp_path / "cours_notes.pptx"
    create_sample_pptx(pptx_path, include_notes=True)

    parser = PptxParser()
    result = parser.parse(pptx_path)

    # Vérifications des notes
    assert "> 🎙️ **Notes de l'orateur :**" in result
    assert "Insister sur la loi de Frank-Starling !" in result
    assert "Question classique d'internat." in result
    assert "Attention : retenir les unités" in result


def test_pptx_parser_table_conversion(tmp_path):
    """Vérifie que les tableaux PowerPoint sont fidèlement traduits en syntaxe Markdown."""
    pptx_path = tmp_path / "cours_table.pptx"
    create_sample_pptx(pptx_path, include_table=True)

    parser = PptxParser()
    result = parser.parse(pptx_path)

    # Vérifications de la syntaxe Markdown du tableau
    assert "| Paramètre" in result
    assert "| Norme" in result
    assert "120/80 mmHg" in result
    assert "60-100 bpm" in result


def test_pptx_parser_image_extraction(tmp_path):
    """Vérifie que les images intégrées sur les diapositives sont archivées dans MediaManager."""
    pptx_path = tmp_path / "cours_image.pptx"
    create_sample_pptx(pptx_path, include_image=True)

    mock_media_mgr = MagicMock()
    mock_media_mgr.store_media_bytes.return_value = MagicMock(filename="stored_fig_123.png")

    parser = PptxParser(media_manager=mock_media_mgr)
    result = parser.parse(pptx_path)

    assert mock_media_mgr.store_media_bytes.called
    assert "stored_fig_123.png" in result
    assert "![Figure" in result


def test_pptx_parser_pagination_and_chunks(tmp_path):
    """Vérifie que chaque diapositive devient une page délimitée et un chunk avec son page_number."""
    pptx_path = tmp_path / "cours_chunks.pptx"
    create_sample_pptx(pptx_path)

    parser = PptxParser()
    raw_md = parser.parse(pptx_path)

    assert "<!-- PAGE: 1 -->" in raw_md
    assert "<!-- PAGE: 2 -->" in raw_md
    assert "[SPLIT]" in raw_md

    chunks = ChunkingService.extract_chunks(raw_md, file_type="pptx")
    assert len(chunks) == 2

    # Chunk 1 : Diapositive 1
    assert chunks[0]["page_number"] == 1
    assert "Cardiaque" in chunks[0]["heading_path"]
    assert "Débit Cardiaque" in chunks[0]["content"]

    # Chunk 2 : Diapositive 2
    assert chunks[1]["page_number"] == 2
    assert "Tableau" in chunks[1]["heading_path"] or "Diapositive 2" in chunks[1]["heading_path"]
    assert "120/80 mmHg" in chunks[1]["content"]


def test_pptx_parser_file_not_found():
    """Vérifie qu'un FileNotFoundError explicite est levé si le fichier est manquant."""
    parser = PptxParser()
    with pytest.raises(FileNotFoundError):
        parser.parse("inexistant.pptx")


def test_pptx_parser_cancellation(tmp_path):
    """Vérifie que le callback check_cancel interrompt proprement l'extraction."""
    pptx_path = tmp_path / "cours_cancel.pptx"
    create_sample_pptx(pptx_path)

    call_count = 0

    def cancel_check() -> bool:
        nonlocal call_count
        call_count += 1
        return call_count > 1

    parser = PptxParser()
    result = parser.parse(pptx_path, check_cancel=cancel_check)

    assert "<!-- PAGE: 1 -->" in result
    assert "<!-- PAGE: 2 -->" not in result


def test_document_parser_pptx_delegation(tmp_path):
    """Vérifie que DocumentParser délègue à PptxParser avec succès."""
    pptx_path = tmp_path / "cours_doc_parser.pptx"
    create_sample_pptx(pptx_path)

    doc_parser = DocumentParser()
    result = doc_parser.parse_document(str(pptx_path))

    assert "## Diapositive 1" in result
    assert "> 🎙️ **Notes de l'orateur :**" in result
    assert "| Paramètre" in result
