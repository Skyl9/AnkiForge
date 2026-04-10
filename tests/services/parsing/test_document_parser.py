import subprocess
import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest
import docx
from pptx import Presentation

from ankiforge.services.parsing.document_parser import DocumentParser


def test_parse_document_file_not_found():
    """Test 1: Si le fichier n'existe pas, ça doit crasher tout de suite."""
    parser = DocumentParser()
    with pytest.raises(FileNotFoundError) as exc_info:
        parser.parse_document("fichier_fantome.pdf")

    assert "est introuvable" in str(exc_info.value)


def test_parse_document_unsupported_format(tmp_path):
    """Test 2: On passe un format exotique non géré (.epub)."""
    fake_epub = tmp_path / "document.epub"
    fake_epub.write_text("fake content")

    parser = DocumentParser()
    with pytest.raises(ValueError) as exc_info:
        parser.parse_document(str(fake_epub))

    assert "Format de fichier non supporté : .epub" in str(exc_info.value)


def test_parse_text_file(tmp_path):
    """Test 3: Vérifie que les fichiers texte sont lus instantanément."""
    fake_md = tmp_path / "cours.md"
    fake_md.write_text("# Titre\nCeci est un cours.", encoding="utf-8")

    parser = DocumentParser()
    mock_callback = MagicMock()
    result = parser.parse_document(str(fake_md), progress_callback=mock_callback)

    assert "Ceci est un cours." in result
    mock_callback.assert_called_with("Lecture du fichier texte immédiate...")


# ==========================================
# TESTS BUREAUTIQUE (DOCX / PPTX)
# ==========================================

def test_parse_docx_success(tmp_path):
    """Test 4: Vérifie l'extraction d'un Word et la traduction des styles en Markdown."""
    fake_docx_path = tmp_path / "test_cours.docx"

    # Création d'un vrai document Word en mémoire
    doc = docx.Document()
    doc.add_heading('Le Théorème de Pythagore', level=1)
    doc.add_paragraph('Voici le contenu du théorème.')
    doc.add_heading('Démonstration', level=2)
    doc.save(fake_docx_path)

    parser = DocumentParser()
    result = parser.parse_document(str(fake_docx_path))

    # Vérifications du Markdown généré
    assert "# Le Théorème de Pythagore" in result
    assert "Voici le contenu du théorème." in result
    assert "## Démonstration" in result


def test_parse_pptx_success(tmp_path):
    """Test 5: Vérifie l'extraction d'un PowerPoint slide par slide avec balise SPLIT."""
    fake_pptx_path = tmp_path / "test_prez.pptx"

    # Création d'une vraie présentation PowerPoint en mémoire
    prs = Presentation()

    # Slide 1
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "Titre Slide 1"
    slide1.placeholders[1].text = "Contenu Slide 1"

    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[0])
    slide2.shapes.title.text = "Titre Slide 2"

    prs.save(fake_pptx_path)

    parser = DocumentParser()
    result = parser.parse_document(str(fake_pptx_path))

    # Vérifications de la structure générée
    assert "## Diapositive 1" in result
    assert "Titre Slide 1" in result
    assert "Contenu Slide 1" in result
    assert "\n\n[SPLIT]\n\n" in result
    assert "## Diapositive 2" in result
    assert "Titre Slide 2" in result


@patch('ankiforge.services.parsing.document_parser.docx', None)
def test_parse_docx_missing_lib(tmp_path):
    """Test 6: Si l'import de python-docx échoue, on doit lever une erreur claire."""
    fake_docx = tmp_path / "test.docx"
    fake_docx.touch()  # Le fichier doit exister pour passer le premier check

    parser = DocumentParser()
    with pytest.raises(RuntimeError) as exc_info:
        parser.parse_document(str(fake_docx))

    assert "python-docx n'est pas installé" in str(exc_info.value)


# ==========================================
# TESTS MARKER (PDF)
# ==========================================

@patch('ankiforge.services.parsing.document_parser.MediaManager')
@patch('subprocess.Popen')
def test_parse_pdf_with_marker_success(mock_popen, MockMediaManager, tmp_path):
    """Test 7: Simule l'extraction d'un PDF avec Marker."""
    fake_pdf = tmp_path / "physique.pdf"
    fake_pdf.write_bytes(b"%PDF-1.4 mock pdf content")

    mock_media_instance = MockMediaManager.return_value
    mock_media_instance.process_extracted_folder.return_value = "Markdown final traité avec images"

    def side_effect_popen(cmd, *args, **kwargs):
        out_dir_idx = cmd.index("--output_dir") + 1
        temp_dir = Path(cmd[out_dir_idx])

        marker_folder = temp_dir / "physique"
        marker_folder.mkdir(parents=True, exist_ok=True)
        (marker_folder / "physique.md").write_text("Contenu brut", encoding="utf-8")

        mock_process = MagicMock()
        mock_process.__enter__.return_value = mock_process
        mock_process.stdout.readline.side_effect = ["Loading AI...\n", "Page 1...\n", ""]
        mock_process.returncode = 0
        return mock_process

    mock_popen.side_effect = side_effect_popen

    parser = DocumentParser()
    mock_callback = MagicMock()
    result = parser.parse_document(str(fake_pdf), progress_callback=mock_callback)

    assert result == "Markdown final traité avec images"
    mock_callback.assert_any_call("Loading AI...")


@patch('subprocess.Popen')
def test_parse_pdf_marker_crash(mock_popen, tmp_path):
    """Test 8: Vérifie ce qui se passe si l'IA plante."""
    fake_pdf = tmp_path / "crash.pdf"
    fake_pdf.write_bytes(b"%PDF")

    def side_effect_popen(*args, **kwargs):
        mock_process = MagicMock()
        mock_process.__enter__.return_value = mock_process
        mock_process.stdout.readline.side_effect = ["Fatal Error", ""]
        mock_process.returncode = 1
        return mock_process

    mock_popen.side_effect = side_effect_popen

    parser = DocumentParser()
    with pytest.raises(RuntimeError) as exc_info:
        parser.parse_document(str(fake_pdf))

    assert "Marker a échoué avec le code erreur 1" in str(exc_info.value)


@patch('subprocess.Popen')
def test_parse_pdf_no_md_generated(mock_popen, tmp_path):
    """Test 9: L'IA dit qu'elle a fini, mais le .md n'est pas là !"""
    fake_pdf = tmp_path / "vide.pdf"
    fake_pdf.write_bytes(b"%PDF")

    def side_effect_popen(*args, **kwargs):
        mock_process = MagicMock()
        mock_process.__enter__.return_value = mock_process
        mock_process.stdout.readline.side_effect = [""]
        mock_process.returncode = 0
        return mock_process

    mock_popen.side_effect = side_effect_popen

    parser = DocumentParser()
    with pytest.raises(FileNotFoundError) as exc_info:
        parser.parse_document(str(fake_pdf))

    assert "Marker n'a pas généré de fichier .md" in str(exc_info.value)


@patch('subprocess.Popen')
def test_parse_pdf_marker_not_installed(mock_popen, tmp_path):
    """Test 10: Si Marker n'est pas sur le PC."""
    fake_pdf = tmp_path / "no_marker.pdf"
    fake_pdf.write_bytes(b"%PDF")

    mock_popen.side_effect = FileNotFoundError("No such file or directory: 'marker_single'")

    parser = DocumentParser()
    with pytest.raises(RuntimeError) as exc_info:
        parser.parse_document(str(fake_pdf))

    assert "Marker n'est pas installé ou introuvable" in str(exc_info.value)