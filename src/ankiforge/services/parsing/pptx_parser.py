"""Service de parsage de présentations PowerPoint (.pptx).

Extrait diapositive par diapositive le texte, les listes à puces, les tableaux,
les images intégrées vers le MediaManager, et surtout les notes secrètes de l'orateur (Speaker Notes),
avec découpage paginé compatible ChunkingService (<!-- PAGE: N -->).
"""

import logging
import mimetypes
from collections.abc import Callable
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ankiforge.services.cards.media_manager import MediaManager

logger = logging.getLogger(__name__)


class PptxParser:
    """Parseur enrichi de diaporamas PowerPoint (.pptx).

    Génère un Markdown structuré par diapositive comprenant :
    - Le titre et le texte hiérarchisé de la diapositive.
    - Les tableaux convertis en syntaxe Markdown (| col | ...).
    - Les illustrations et schémas archivés dans le MediaManager.
    - Les notes de présentation de l'orateur (Speaker Notes) sous forme de citation enrichie.
    """

    def __init__(self, media_manager: MediaManager | None = None) -> None:
        """Initialise le parseur avec un gestionnaire de médias pour l'archivage d'images."""
        self.media_manager = media_manager or MediaManager()

    def parse(
        self,
        pptx_path: str | Path,
        progress_callback: Callable[[str], None] | None = None,
        check_cancel: Callable[[], bool] | None = None,
    ) -> str:
        """Parse l'ensemble du fichier PowerPoint et retourne le Markdown paginé.

        Args:
            pptx_path: Chemin vers le fichier .pptx.
            progress_callback: Callback optionnel pour signaler l'avancement.
            check_cancel: Callback optionnel pour vérifier l'annulation utilisateur.

        Returns:
            str: Le texte Markdown complet avec délimiteurs <!-- PAGE: N --> et [SPLIT].

        Raises:
            FileNotFoundError: Si le fichier PPTX n'existe pas.
            RuntimeError: En cas d'échec de lecture du fichier.
        """
        path_obj = Path(pptx_path)
        if not path_obj.exists():
            raise FileNotFoundError(f"Le fichier PowerPoint '{path_obj}' est introuvable.")

        logger.info("Début du parsage enrichi PPTX : %s", path_obj.name)

        try:
            prs = Presentation(str(path_obj))
        except Exception as err:
            logger.error("Impossible d'ouvrir la présentation PPTX '%s' : %s", path_obj.name, err)
            raise RuntimeError(f"Erreur lors de la lecture du fichier PowerPoint : {err}") from err

        total_slides = len(prs.slides)
        slide_outputs: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            if check_cancel and check_cancel():
                logger.warning("Extraction PPTX interrompue par l'utilisateur à la diapositive %d.", idx)
                break

            if progress_callback:
                progress_callback(f"Analyse de la diapositive {idx}/{total_slides}...")

            slide_md = self._process_slide(slide=slide, slide_idx=idx)
            if slide_md.strip():
                chunk_header = f"<!-- PAGE: {idx} -->"
                slide_outputs.append(f"{chunk_header}\n\n{slide_md.strip()}")

        if not slide_outputs:
            logger.warning("Aucune diapositive exploitable dans le PowerPoint : %s", path_obj.name)
            return f"# {path_obj.stem}\n\n*Présentation vide ou sans contenu textuel.*"

        logger.info(
            "Extraction PPTX achevée avec succès : %d diapositives extraites pour '%s'",
            len(slide_outputs),
            path_obj.name,
        )
        return "\n\n[SPLIT]\n\n".join(slide_outputs)

    def _process_slide(self, slide: Any, slide_idx: int) -> str:
        """Extrait le contenu complet d'une diapositive (titre, formes, tableaux, images, notes)."""
        # 1. Extraction du titre officiel de la diapositive
        title_shape = getattr(slide.shapes, "title", None)
        slide_title = ""
        if title_shape and hasattr(title_shape, "text") and title_shape.text.strip():
            slide_title = title_shape.text.strip().replace("\n", " ")

        # 2. Extraction du corps de la diapositive (texte, tableaux, images)
        body_elements: list[str] = []
        image_counter = 0

        def extract_shapes(shapes: Any) -> None:
            nonlocal image_counter
            for shape in shapes:
                # Éviter de dupliquer le titre déjà extrait
                if shape == title_shape:
                    continue

                # A. Groupes de formes
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.GROUP or hasattr(shape, "shapes"):
                    extract_shapes(shape.shapes)
                    continue

                # B. Tableaux
                if getattr(shape, "has_table", False):
                    table_md = self._format_table_markdown(shape.table)
                    if table_md:
                        body_elements.append(table_md)
                    continue

                # C. Images / Schémas
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                    try:
                        image = shape.image
                        image_bytes = image.blob
                        ext = getattr(image, "ext", "png").lower()
                        image_counter += 1
                        orig_name = f"slide_{slide_idx}_fig_{image_counter}.{ext}"
                        mime_type, _ = mimetypes.guess_type(orig_name)

                        media = self.media_manager.store_media_bytes(
                            data=image_bytes,
                            original_name=orig_name,
                            mime_type=mime_type or f"image/{ext}",
                        )
                        if media:
                            body_elements.append(f"![Figure {image_counter} (Diapositive {slide_idx})]({media.filename})")
                    except Exception as err:
                        logger.debug("Impossible d'extraire une image sur la diapositive %d : %s", slide_idx, err)
                    continue

                # D. Cadres de texte
                if getattr(shape, "has_text_frame", False) and shape.text_frame:
                    text_blocks: list[str] = []
                    for para in shape.text_frame.paragraphs:
                        p_text = para.text.strip()
                        if not p_text:
                            continue
                        level = getattr(para, "level", 0) or 0
                        if level > 0:
                            indent = "  " * level
                            text_blocks.append(f"{indent}- {p_text}")
                        else:
                            text_blocks.append(p_text)

                    if text_blocks:
                        body_elements.append("\n".join(text_blocks))

        extract_shapes(slide.shapes)

        # 3. Extraction des notes de l'orateur (Speaker Notes)
        speaker_notes_md = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes_slide = slide.notes_slide
                if notes_slide and hasattr(notes_slide, "notes_text_frame"):
                    raw_notes = notes_slide.notes_text_frame.text.strip()
                    if raw_notes:
                        # Mettre les notes sous forme de citation Markdown bien délimitée
                        formatted_notes = "\n> ".join(raw_notes.splitlines())
                        speaker_notes_md = f"> 🎙️ **Notes de l'orateur :**\n> {formatted_notes}"
            except Exception as err:
                logger.debug("Impossible d'extraire les notes de l'orateur pour la slide %d : %s", slide_idx, err)

        # 4. Assemblage final de la diapositive
        header_title = f"Diapositive {slide_idx}"
        if slide_title:
            header_title = f"{header_title} : {slide_title}"

        parts = [f"## {header_title}"]
        if body_elements:
            parts.append("\n\n".join(body_elements))
        if speaker_notes_md:
            parts.append(speaker_notes_md)

        return "\n\n".join(parts)

    @staticmethod
    def _format_table_markdown(table: Any) -> str:
        """Convertit un tableau python-pptx en syntaxe de tableau Markdown."""
        rows_data: list[list[str]] = []
        for row in table.rows:
            row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows_data.append(row_cells)

        if not rows_data or not any(any(cell for cell in r) for r in rows_data):
            return ""

        num_cols = max(len(r) for r in rows_data)
        # Normaliser les longueurs des lignes
        for r in rows_data:
            while len(r) < num_cols:
                r.append("")

        headers = rows_data[0]
        # Si la première ligne est vide, donner des en-têtes par défaut
        if not any(headers):
            headers = [f"Colonne {i + 1}" for i in range(num_cols)]

        col_widths = [max(len(str(r[c])) for r in rows_data) for c in range(num_cols)]
        col_widths = [max(w, len(headers[c]), 3) for c, w in enumerate(col_widths)]

        header_line = "| " + " | ".join(headers[c].ljust(col_widths[c]) for c in range(num_cols)) + " |"
        sep_line = "| " + " | ".join("-" * col_widths[c] for c in range(num_cols)) + " |"

        body_lines: list[str] = []
        for row in rows_data[1:]:
            row_str = "| " + " | ".join(row[c].ljust(col_widths[c]) for c in range(num_cols)) + " |"
            body_lines.append(row_str)

        return "\n".join([header_line, sep_line] + body_lines)
